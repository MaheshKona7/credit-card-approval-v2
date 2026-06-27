from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Define color scheme (Fintech theme: Indigo/Navy, Slate, White)
    COLOR_TITLE_BG = RGBColor(17, 24, 39)     # Dark slate
    COLOR_PRIMARY = RGBColor(99, 102, 241)     # Indigo blue
    COLOR_TEXT = RGBColor(55, 65, 81)          # Dark grey
    COLOR_WHITE = RGBColor(255, 255, 255)
    
    # ------------------ SLIDE 1: Title Slide ------------------
    slide_layout = prs.slide_layouts[0] # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Credit Card Approval Prediction"
    subtitle.text = "Automated Underwriting & Machine Learning Dashboard\nJune 2026"
    
    # Format Title Slide text colors
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # ------------------ SLIDE 2: Project Overview ------------------
    slide_layout = prs.slide_layouts[1] # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = "1. Project Overview & Business Value"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Automated Credit Card Decisioning Platform"
    
    p = tf.add_paragraph()
    p.text = "• Underwriting Automation: Replaces manual paper auditing with an instant predictive decisioning model."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Credit Risk Mitigation: Screens applicants for high-risk repayment histories, identifying profiles that are over 30 days past due."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Scalable Platform: Integrates data visualization, preprocessing, machine learning model pipeline, and an interactive Flask dashboard."
    p.level = 1
    
    # ------------------ SLIDE 3: Features & Engineering ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "2. Underwriting Features & Engineering"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Predictive Features Captured on the Underwriting Form:"
    
    p = tf.add_paragraph()
    p.text = "• Demographic Profiles: Gender, Car Owner, Realty Owner, Household size."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Financial Metrics: Total Annual Income, Income Source Category (Working, Commercial associate, Pensioner, State servant, Student)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Employment & Stability: Occupation type (19 sectors), Age (DAYS_BIRTH), and Employment Duration (DAYS_EMPLOYED) converted from years to negative days."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Target Mappings: High-risk delinquents mapped to label '1' (Rejected) if they are 30+ days past due. Good accounts mapped to '0' (Approved)."
    p.level = 1

    # ------------------ SLIDE 4: Visual Analytics (EDA) ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "3. Exploratory Data Analysis Insights"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Generated Visualization Insights (Saved inside static/images/):"
    
    p = tf.add_paragraph()
    p.text = "• Approval Distribution Count: Identifies overall risk proportion (88% approved, 12% rejected due to delinquent records)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Annual Income KDE Density: Compares income scale density between approved and rejected accounts."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Education & Income Source Plots: Stacked charts showing success/failure rates segmented by academic credentials and employment type."
    p.level = 1
    
    # ------------------ SLIDE 5: Machine Learning Results ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "4. Machine Learning Model Evaluation"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Algorithms Evaluation on 36,457 integrated applicant profiles:"
    
    p = tf.add_paragraph()
    p.text = "• Random Forest Classifier (Production Model):"
    p.level = 1
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "  - Accuracy: 88.95%  |  Precision (C1): 54.87%  |  Recall (C1): 34.15%  |  F1-Score: 42.10%"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Decision Tree Classifier: Accuracy 88.30%  |  F1-Score (C1): 39.29%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Gradient Boosting (fallback): Accuracy 88.25%  |  F1-Score (C1): 0.70%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Logistic Regression: Accuracy 88.23%  |  F1-Score (C1): 0.00% (predicts majority class)"
    p.level = 1
    
    # ------------------ SLIDE 6: Flask Dashboard ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "5. Underwriting Dashboard (Flask App)"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Flask Web Application Architecture & Styling:"
    
    p = tf.add_paragraph()
    p.text = "• Main Landing Dashboard: Highlights project metadata, model statistics cards, and embeds custom EDA charts."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Underwriting Form: Fully interactive form taking user inputs. Categorical values are pre-coded to match Random Forest parameters."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Decision View: Dynamically renders credit decisions with color-coded outcome badges (Green for APPROVED, Red for REJECTED)."
    p.level = 1
    
    # ------------------ SLIDE 7: How to Run ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "6. System Operation Guide"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    tf = slide.placeholders[1].text_frame
    tf.text = "How to run the pipeline and server locally:"
    
    p = tf.add_paragraph()
    p.text = "1. Navigate to directory and activate virtual environment:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   PowerShell:  .\\venv\\Scripts\\Activate.ps1"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "2. Generate data visualizations (EDA):"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   Command:     python eda.py"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "3. Train models & export production pickle:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   Command:     python model.py"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "4. Host the interactive Flask application:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   Command:     python app.py  (Access at http://127.0.0.1:5000/)"
    p.level = 2
    
    # Save presentation
    prs.save("Credit_Card_Approval_System_Presentation.pptx")
    print("PowerPoint presentation generated successfully as Credit_Card_Approval_System_Presentation.pptx!")

if __name__ == "__main__":
    create_presentation()
